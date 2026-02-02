"""
Phase D: Build PPTX from deck_json and generated images.

Supports three layouts:
  - TITLE              — Cover slide (title + subtitle)
  - TITLE_LEFT_IMAGE_RIGHT — Left text + right image (default)
  - IMAGE_FULL_TEXT_BOTTOM  — Full-bleed image + text strip at bottom

Each content slide renders:
  - title   (in a colored header bar)
  - message (key takeaway line, bold)
  - body    (80-180 char description paragraph)
  - bullets (3-5 items)
  - image   (Imagen result or fallback shape)

If no image file is available, a PPT-native shape diagram is drawn
directly on the slide so that no slide is ever blank.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
try:
    from PIL import Image
except ImportError:
    Image = None

def _hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(43, 87, 154)
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

def _set_slide_bg(slide, rgb_color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def _create_left_column(slide, title, takeaway, bullets, body, font_name):
    """
    Creates the left column content (60% width):
    Title -> Takeaway -> Bullets -> Body
    """
    LEFT_MARGIN = Inches(0.5)
    TOP_MARGIN = Inches(0.5)
    COL_WIDTH = Inches(7.5)  # Approx 60% of 13.33 inches

    current_top = TOP_MARGIN

    # 1. Title (Heading)
    if title:
        tx = slide.shapes.add_textbox(LEFT_MARGIN, current_top, COL_WIDTH, Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        current_top += Inches(1.3)

    # 2. Takeaway (Emphasis Box)
    if takeaway:
        box_h = Inches(1.1)
        # Background box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, current_top, COL_WIDTH, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 240, 255)  # Light Blue bg
        shape.line.color.rgb = RGBColor(100, 150, 230)
        
        # Text inside
        tx = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), current_top + Inches(0.1), COL_WIDTH - Inches(0.4), box_h - Inches(0.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = takeaway
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0, 50, 100)
        p.alignment = PP_ALIGN.LEFT
        
        current_top += box_h + Inches(0.3)

    # 3. Bullets
    if bullets:
        # Estimated height for 3 bullets
        bullets_h = Inches(2.0)
        tx = slide.shapes.add_textbox(LEFT_MARGIN, current_top, COL_WIDTH, bullets_h)
        tf = tx.text_frame
        tf.word_wrap = True
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"\u2022 {b}"
            p.font.size = Pt(24)
            p.font.name = font_name
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.space_after = Pt(12)
        
        current_top += bullets_h + Inches(0.2)

    # 4. Body (Smaller font, explanatory text)
    if body:
        body_h = Inches(2.5)  # Remaining space roughly
        tx = slide.shapes.add_textbox(LEFT_MARGIN, current_top, COL_WIDTH, body_h)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body
        p.font.size = Pt(20)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.line_spacing = 1.3

def _add_image_contain(slide, image_path, left, top, max_width, max_height):
    """
    Add image fitting effectively within the box (Contain) maintaining aspect ratio.
    scale = min(box_w/img_w, box_h/img_h)
    """
    if not image_path or not os.path.exists(image_path):
        return

    # Get image dimensions
    img_w, img_h = 100, 100 # defaults
    if Image:
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
        except Exception:
            pass
    
    # Calculate scale
    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)
    
    final_w = img_w * scale
    final_h = img_h * scale
    
    # Center it in the box
    final_left = left + (max_width - final_w) / 2
    final_top = top + (max_height - final_h) / 2
    
    slide.shapes.add_picture(
        image_path, final_left, final_top, width=final_w, height=final_h
    )

def generate_pptx(plan, image_paths, output_path="presentation.pptx", title="Presentation"):
    """
    Generates PPTX for Video Presentation (2-column layout).
    """
    prs = Presentation()
    # Widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Extract theme
    theme = plan.get('theme', {})
    font_name = theme.get('font', 'Meiryo')
    primary = _hex_to_rgb(theme.get('primary_color', '#2B579A'))

    # ── Cover Slide ──
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(cover, primary)
    
    # Cover title
    tx = cover.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.name = font_name
    p.alignment = PP_ALIGN.CENTER

    # ── Content Slides ──
    # ── Content Slides ──
    for i, sd in enumerate(plan['slides']):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, RGBColor(255, 255, 255)) # White bg for content slides
        
        # 1. Left Content (Title -> Takeaway -> Bullets -> Body)
        _create_left_column(
            slide,
            sd.get('title', ''),
            sd.get('takeaway', ''),
            sd.get('bullets', []),
            sd.get('body', ''),
            font_name
        )
        
        # 2. Right Content (Image) - 40% area
        # Slide Width 13.33. 
        # Left 60% is roughly 8.0.
        # Right area: Left=8.2, Top=0.5, Width=4.6, Height=6.5
        img_left = Inches(8.2)
        img_top = Inches(0.5)
        img_w = Inches(4.6)
        img_h = Inches(6.5)
        
        # Visual Frame for Image Area (Optional, but good for layout debugging/structure)
        # bg_box = slide.shapes.add_shape(
        #     MSO_SHAPE.RECTANGLE, img_left, img_top, img_w, img_h
        # )
        # bg_box.fill.solid()
        # bg_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        # bg_box.line.color.rgb = RGBColor(230, 230, 230)
        
        img_path = image_paths.get(i)
        if img_path:
            _add_image_contain(slide, img_path, img_left, img_top, img_w, img_h)

    prs.save(output_path)
    return output_path
